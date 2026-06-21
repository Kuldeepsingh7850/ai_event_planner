"""
Generate AI Event Planner Project PPT
Uses python-pptx to create a professional presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
CARD_BG       = RGBColor(0x1A, 0x23, 0x3B)   # Slightly lighter navy
ACCENT_BLUE   = RGBColor(0x38, 0xBD, 0xF8)   # Bright cyan-blue
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)   # Soft purple
ACCENT_GREEN  = RGBColor(0x4A, 0xDE, 0x80)   # Green
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)   # Amber/Orange
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCB, 0xD5, 0xE1)   # Slate-300
MID_GRAY      = RGBColor(0x94, 0xA3, 0xB8)   # Slate-400
GRADIENT_START = RGBColor(0x1E, 0x29, 0x3B)
TABLE_HEADER  = RGBColor(0x1E, 0x40, 0x6E)
TABLE_ROW_ALT = RGBColor(0x16, 0x1E, 0x33)
TABLE_ROW     = RGBColor(0x11, 0x19, 0x2B)
BORDER_COLOR  = RGBColor(0x33, 0x44, 0x5C)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helpers ─────────────────────────────────────────────────────

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust corner radius
    if radius is not None:
        shape.adjustments[0] = radius
    else:
        shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, font_size=16, color=WHITE, bold=False, line_spacing=1.5, font_name="Calibri", alignment=PP_ALIGN.LEFT):
    """lines is a list of (text, color, bold, font_size) tuples or just strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.color.rgb = line[1] if len(line) > 1 else color
            p.font.bold = line[2] if len(line) > 2 else bold
            p.font.size = Pt(line[3] if len(line) > 3 else font_size)
        else:
            p.text = line
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.size = Pt(font_size)
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)
    return txBox

def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_slide_number(slide, num, total=13):
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(0.9), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_bottom_bar(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.35), Inches(13.333), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

def add_section_title(slide, title, subtitle=None, slide_num=1):
    set_slide_bg(slide, DARK_BG)
    add_bottom_bar(slide)
    # Title accent line
    add_accent_line(slide, Inches(0.8), Inches(1.4), Inches(1.2), ACCENT_BLUE)
    add_text_box(slide, Inches(0.8), Inches(1.55), Inches(11), Inches(0.8),
                 title, font_size=36, color=WHITE, bold=True, font_name="Calibri")
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(2.3), Inches(11), Inches(0.5),
                     subtitle, font_size=18, color=MID_GRAY, font_name="Calibri")
    add_slide_number(slide, slide_num)

def make_table(slide, left, top, width, rows_data, col_widths, header_color=TABLE_HEADER):
    """rows_data: list of lists, first row is header"""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.45 * num_rows))
    table = table_shape.table

    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if ri == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.alignment = PP_ALIGN.LEFT
            # Cell fill
            cf = cell.fill
            cf.solid()
            if ri == 0:
                cf.fore_color.rgb = header_color
            elif ri % 2 == 0:
                cf.fore_color.rgb = TABLE_ROW_ALT
            else:
                cf.fore_color.rgb = TABLE_ROW
            # Borders
            cell.margin_left = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
    return table_shape


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK_BG)

# Decorative top gradient bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Main card
card = add_shape_rect(slide, Inches(1.5), Inches(1.0), Inches(10.333), Inches(5.2), CARD_BG, BORDER_COLOR, 0.03)

# "MAJOR PROJECT PRESENTATION ON" subtitle
add_text_box(slide, Inches(2.0), Inches(1.5), Inches(9.333), Inches(0.5),
             "MAJOR PROJECT PRESENTATION ON", font_size=14, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Accent line under subtitle
add_accent_line(slide, Inches(5.5), Inches(2.05), Inches(2.333), ACCENT_BLUE)

# Project Title
add_text_box(slide, Inches(2.0), Inches(2.3), Inches(9.333), Inches(1.0),
             "JAGAH — AI Event Planner", font_size=44, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Tagline
add_text_box(slide, Inches(2.0), Inches(3.3), Inches(9.333), Inches(0.6),
             "An Intelligent Full-Stack Event Sourcing & Organization Ecosystem",
             font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Separator
add_accent_line(slide, Inches(4.5), Inches(4.1), Inches(4.333), ACCENT_PURPLE)

# Presented By
add_text_box(slide, Inches(2.0), Inches(4.4), Inches(9.333), Inches(0.4),
             "Presented By", font_size=14, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_text_box(slide, Inches(2.0), Inches(4.8), Inches(9.333), Inches(0.5),
             "Mangal Sharma", font_size=28, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Bottom bar
add_bottom_bar(slide)
add_slide_number(slide, 1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTENTS
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "CONTENTS", None, 2)

contents = [
    ("01", "Introduction"),
    ("02", "Objectives"),
    ("03", "Tech Stack"),
    ("04", "Functionalities"),
    ("05", "Software Requirements"),
    ("06", "Hardware Requirements"),
    ("07", "Database Design"),
    ("08", "System Workflows"),
    ("09", "Snapshots"),
    ("10", "Conclusion"),
]

# Two columns layout
for i, (num, title) in enumerate(contents):
    col = 0 if i < 5 else 1
    row = i if i < 5 else i - 5
    x = Inches(1.2 + col * 5.5)
    y = Inches(3.0 + row * 0.75)

    # Number badge
    badge = add_shape_rect(slide, x, y, Inches(0.5), Inches(0.45), ACCENT_BLUE, radius=0.15)
    badge.text_frame.paragraphs[0].text = num
    badge.text_frame.paragraphs[0].font.size = Pt(14)
    badge.text_frame.paragraphs[0].font.bold = True
    badge.text_frame.paragraphs[0].font.color.rgb = DARK_BG
    badge.text_frame.paragraphs[0].font.name = "Calibri"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge.text_frame.word_wrap = False

    add_text_box(slide, x + Inches(0.65), y + Inches(0.03), Inches(4.5), Inches(0.45),
                 title, font_size=20, color=LIGHT_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — INTRODUCTION
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "INTRODUCTION", "Overview of the AI Event Planner System", 3)

intro_points = [
    "JAGAH is an AI-powered event management web application that helps users plan and organize events like weddings, birthdays, farewells, and corporate seminars.",
    "The platform consolidates guest lists, budget tracking, task checklists, vendor management, and AI-generated recommendations into a single unified dashboard.",
    "It integrates Google Gemini, Groq, and Grok APIs to automatically generate custom themes, decoration ideas, food menus, event timelines, and localized venue suggestions.",
    "Features a dual-mode database engine — if MySQL is unavailable, the system automatically switches to an in-memory mock database to keep the application running.",
    "Built using Next.js, Node.js, Express.js, MySQL, and Tailwind CSS v4 with full dark/light theme support.",
]

for i, pt in enumerate(intro_points):
    y = Inches(3.0 + i * 0.85)
    # Bullet dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.08), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_BLUE
    dot.line.fill.background()
    add_text_box(slide, Inches(1.5), y - Inches(0.05), Inches(10.5), Inches(0.85),
                 pt, font_size=16, color=LIGHT_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4 — OBJECTIVES
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "OBJECTIVES", "Primary & Secondary Goals of the Project", 4)

objectives = [
    ("Centralize Event Planning", "Combine budgets, guest RSVPs, checklists, and vendor contacts into one interactive dashboard.", ACCENT_BLUE),
    ("AI Recommendation Engine", "Generate tailored themes, menus, schedules, budget allocations, and venue suggestions using Generative AI.", ACCENT_PURPLE),
    ("Interactive AI Chat Assistant", "Allow users to ask event planning questions and receive contextual advice instantly.", ACCENT_GREEN),
    ("Secure Authentication", "Implement JWT-based sessions with bcrypt hashing and Google OAuth sign-in.", ACCENT_ORANGE),
    ("Database Resilience", "Build an in-memory SQL mock engine so the app runs without MySQL setup.", ACCENT_BLUE),
    ("Admin Moderation Panel", "Provide admins with tools to manage users, audit events, and moderate platform activity.", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(3.0 + row * 2.1)

    card = add_shape_rect(slide, x, y, Inches(3.7), Inches(1.8), CARD_BG, BORDER_COLOR, 0.05)
    # Color accent bar on top of card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.15), y + Inches(0.1), Inches(0.6), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text_box(slide, x + Inches(0.15), y + Inches(0.25), Inches(3.4), Inches(0.4),
                 title, font_size=16, color=color, bold=True, font_name="Calibri")
    add_text_box(slide, x + Inches(0.15), y + Inches(0.7), Inches(3.4), Inches(1.0),
                 desc, font_size=13, color=MID_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "TECH STACK", "Technologies & Frameworks Used", 5)

tech_data = [
    ["Layer", "Technology", "Role"],
    ["Frontend", "Next.js (React)", "Client-side SPA with Pages Router"],
    ["Styling", "Tailwind CSS v4", "Responsive layouts, dark/light themes"],
    ["Backend", "Node.js + Express.js", "REST API server, middleware, routing"],
    ["Database", "MySQL 8.0", "Relational data storage with InnoDB"],
    ["AI Integration", "Google Gemini, Groq, Grok", "Event suggestion generation & chat"],
    ["Authentication", "JWT + bcryptjs", "Stateless session tokens & hashing"],
    ["DB Fallback", "In-Memory Mock Engine", "Simulated SQL without MySQL"],
]

make_table(slide, Inches(1.0), Inches(3.0), Inches(11.333), tech_data,
           [Inches(2.5), Inches(3.5), Inches(5.333)])


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6 — FUNCTIONALITIES
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "FUNCTIONALITIES", "User & Admin Feature Set", 6)

# User features card
user_card = add_shape_rect(slide, Inches(0.8), Inches(2.8), Inches(5.8), Inches(4.2), CARD_BG, BORDER_COLOR, 0.03)
add_text_box(slide, Inches(1.1), Inches(2.9), Inches(3.0), Inches(0.4),
             "👤  User Features", font_size=20, color=ACCENT_BLUE, bold=True)
add_accent_line(slide, Inches(1.1), Inches(3.35), Inches(2.0), ACCENT_BLUE)

user_features = [
    "Register / Login (Local + Google OAuth)",
    "Create, update, and delete events",
    "Add and track guest RSVPs (Going / Pending / Not Going)",
    "Log expenses with auto budget recalculation",
    "Manage task checklists with deadlines",
    "Generate AI suggestions (themes, menus, venues)",
    "Chat with AI planning assistant",
]
for i, feat in enumerate(user_features):
    y = Inches(3.55 + i * 0.45)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), y + Inches(0.08), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_BLUE
    dot.line.fill.background()
    add_text_box(slide, Inches(1.55), y - Inches(0.02), Inches(4.8), Inches(0.45),
                 feat, font_size=14, color=LIGHT_GRAY, font_name="Calibri")

# Admin features card
admin_card = add_shape_rect(slide, Inches(6.9), Inches(2.8), Inches(5.6), Inches(4.2), CARD_BG, BORDER_COLOR, 0.03)
add_text_box(slide, Inches(7.2), Inches(2.9), Inches(3.0), Inches(0.4),
             "🛡️  Admin Features", font_size=20, color=ACCENT_PURPLE, bold=True)
add_accent_line(slide, Inches(7.2), Inches(3.35), Inches(2.0), ACCENT_PURPLE)

admin_features = [
    "View all registered users and events",
    "Block / Unblock user accounts",
    "Change user roles (User ↔ Admin)",
    "Delete user accounts",
    "Monitor platform-wide statistics",
]
for i, feat in enumerate(admin_features):
    y = Inches(3.55 + i * 0.45)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.4), y + Inches(0.08), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_PURPLE
    dot.line.fill.background()
    add_text_box(slide, Inches(7.65), y - Inches(0.02), Inches(4.5), Inches(0.45),
                 feat, font_size=14, color=LIGHT_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7 — SOFTWARE REQUIREMENTS
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "SOFTWARE REQUIREMENTS", "Minimum Software Specifications", 7)

sw_data = [
    ["Component", "Specification"],
    ["Operating System", "Windows 10/11, macOS, or Linux (Ubuntu 20.04+)"],
    ["Runtime", "Node.js (v18.0.0+) & NPM (v9.0.0+)"],
    ["Database", "MySQL 8.0+ or MariaDB 10.5+"],
    ["IDE", "Visual Studio Code"],
    ["API Testing", "Postman or Thunder Client"],
    ["Browser", "Google Chrome, Mozilla Firefox, or Microsoft Edge"],
]

make_table(slide, Inches(1.5), Inches(3.0), Inches(10.333), sw_data,
           [Inches(3.5), Inches(6.833)])


# ═══════════════════════════════════════════════════════════════════
# SLIDE 8 — HARDWARE REQUIREMENTS
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "HARDWARE REQUIREMENTS", "Development & Deployment Specifications", 8)

hw_data = [
    ["Feature", "Development Environment", "Server / Deployment"],
    ["Processor", "Intel Core i5 / AMD Ryzen 5 (4 Cores)", "1 vCPU"],
    ["RAM", "8 GB DDR4", "1 GB"],
    ["Storage", "10 GB free SSD", "20 GB SSD"],
    ["Network", "Localhost", "Broadband with Static Public IP"],
]

make_table(slide, Inches(1.5), Inches(3.0), Inches(10.333), hw_data,
           [Inches(2.5), Inches(4.5), Inches(3.333)])


# ═══════════════════════════════════════════════════════════════════
# SLIDE 9 — DATABASE DESIGN
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "DATABASE DESIGN", "MySQL Schema & Relationships", 9)

# Left side — Table list
db_card = add_shape_rect(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(3.8), CARD_BG, BORDER_COLOR, 0.03)
add_text_box(slide, Inches(1.1), Inches(3.1), Inches(3.0), Inches(0.35),
             "9 Database Tables", font_size=18, color=ACCENT_BLUE, bold=True)

tables = ["users", "events", "guests", "budget", "expenses", "tasks", "vendors", "notifications", "feedback"]
for i, t in enumerate(tables):
    col = i % 3
    row = i // 3
    x = Inches(1.2 + col * 1.7)
    y = Inches(3.6 + row * 0.55)
    badge = add_shape_rect(slide, x, y, Inches(1.5), Inches(0.4), GRADIENT_START, BORDER_COLOR, 0.15)
    badge.text_frame.paragraphs[0].text = t
    badge.text_frame.paragraphs[0].font.size = Pt(12)
    badge.text_frame.paragraphs[0].font.color.rgb = ACCENT_BLUE
    badge.text_frame.paragraphs[0].font.name = "Consolas"
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Relationships below tables
add_text_box(slide, Inches(1.1), Inches(5.3), Inches(5.0), Inches(0.3),
             "Key Relationships:", font_size=14, color=ACCENT_PURPLE, bold=True)
rels = [
    "Users → Events (1 : N)",
    "Events → Budget (1 : 1)",
    "Events → Guests, Expenses, Tasks, Vendors (1 : N)",
]
for i, r in enumerate(rels):
    add_text_box(slide, Inches(1.3), Inches(5.65 + i * 0.35), Inches(4.8), Inches(0.35),
                 f"•  {r}", font_size=12, color=LIGHT_GRAY)

# Right side — Constraints
rt_card = add_shape_rect(slide, Inches(6.7), Inches(3.0), Inches(5.8), Inches(3.8), CARD_BG, BORDER_COLOR, 0.03)
add_text_box(slide, Inches(7.0), Inches(3.1), Inches(3.0), Inches(0.35),
             "Design Constraints", font_size=18, color=ACCENT_GREEN, bold=True)

constraints = [
    "ON DELETE CASCADE on all foreign keys",
    "Unique index on users(email)",
    "Unique index on budget(event_id)",
    "Auto-increment primary keys (INT)",
    "Timestamps: created_at with DEFAULT CURRENT_TIMESTAMP",
    "Normalized to 1NF, 2NF, and 3NF",
    "InnoDB storage engine for ACID compliance",
]
for i, c in enumerate(constraints):
    y = Inches(3.6 + i * 0.45)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), y + Inches(0.08), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_GREEN
    dot.line.fill.background()
    add_text_box(slide, Inches(7.45), y - Inches(0.02), Inches(4.8), Inches(0.45),
                 c, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 10 — SYSTEM WORKFLOWS
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "SYSTEM WORKFLOWS", "User & Admin Flow Diagrams", 10)

# User Workflow Card
uw_card = add_shape_rect(slide, Inches(0.8), Inches(3.0), Inches(5.8), Inches(3.8), CARD_BG, BORDER_COLOR, 0.03)
add_text_box(slide, Inches(1.1), Inches(3.1), Inches(3.0), Inches(0.35),
             "👤  User Workflow", font_size=18, color=ACCENT_BLUE, bold=True)
add_accent_line(slide, Inches(1.1), Inches(3.5), Inches(1.5), ACCENT_BLUE)

user_steps = [
    "Registration / Login",
    "Create New Event",
    "Generate AI Suggestions",
    "Manage Guests & RSVPs",
    "Track Budget & Expenses",
    "Task Checklist Management",
    "Chat with AI Assistant",
    "Submit Feedback",
]
for i, step in enumerate(user_steps):
    y = Inches(3.7 + i * 0.40)
    # Step number
    num_shape = add_shape_rect(slide, Inches(1.2), y, Inches(0.3), Inches(0.3), ACCENT_BLUE, radius=0.5)
    num_shape.text_frame.paragraphs[0].text = str(i + 1)
    num_shape.text_frame.paragraphs[0].font.size = Pt(10)
    num_shape.text_frame.paragraphs[0].font.bold = True
    num_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BG
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Arrow connector (except last)
    if i < len(user_steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.33), y + Inches(0.3), Pt(2), Inches(0.1))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_BLUE
        arrow.line.fill.background()
    add_text_box(slide, Inches(1.65), y - Inches(0.02), Inches(4.5), Inches(0.35),
                 step, font_size=13, color=LIGHT_GRAY)

# Admin Workflow Card
aw_card = add_shape_rect(slide, Inches(6.9), Inches(3.0), Inches(5.6), Inches(3.8), CARD_BG, BORDER_COLOR, 0.03)
add_text_box(slide, Inches(7.2), Inches(3.1), Inches(3.0), Inches(0.35),
             "🛡️  Admin Workflow", font_size=18, color=ACCENT_PURPLE, bold=True)
add_accent_line(slide, Inches(7.2), Inches(3.5), Inches(1.5), ACCENT_PURPLE)

admin_steps = [
    "Admin Login",
    "Admin Dashboard Overview",
    "View All Users & Events",
    "User Moderation (Block / Unblock)",
    "Role Management (User ↔ Admin)",
    "Delete User Accounts",
    "Platform Analytics & Audit",
]
for i, step in enumerate(admin_steps):
    y = Inches(3.7 + i * 0.43)
    num_shape = add_shape_rect(slide, Inches(7.4), y, Inches(0.3), Inches(0.3), ACCENT_PURPLE, radius=0.5)
    num_shape.text_frame.paragraphs[0].text = str(i + 1)
    num_shape.text_frame.paragraphs[0].font.size = Pt(10)
    num_shape.text_frame.paragraphs[0].font.bold = True
    num_shape.text_frame.paragraphs[0].font.color.rgb = DARK_BG
    num_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    if i < len(admin_steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.53), y + Inches(0.3), Pt(2), Inches(0.13))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_PURPLE
        arrow.line.fill.background()
    add_text_box(slide, Inches(7.85), y - Inches(0.02), Inches(4.3), Inches(0.35),
                 step, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 11 — SNAPSHOTS
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "SNAPSHOTS", "Application Screenshots & UI Highlights", 11)

# Add screenshots if available
screenshot_dir = r"C:\Users\H\.gemini\antigravity\brain\81ff7180-48d8-40b1-9fcc-8912ee6767ae"
screenshots_added = False

# Try to find and add screenshots
possible_screenshots = []
if os.path.exists(screenshot_dir):
    for f in os.listdir(screenshot_dir):
        if f.lower().endswith(('.png', '.jpg', '.jpeg')) and ('screenshot' in f.lower() or 'guest' in f.lower() or 'dashboard' in f.lower() or 'workflow' in f.lower() or 'db_' in f.lower()):
            possible_screenshots.append(os.path.join(screenshot_dir, f))

if possible_screenshots:
    for i, img_path in enumerate(possible_screenshots[:4]):  # Max 4 screenshots
        col = i % 2
        row = i // 2
        x = Inches(1.0 + col * 6.0)
        y = Inches(3.0 + row * 2.2)
        try:
            slide.shapes.add_picture(img_path, x, y, Inches(5.5), Inches(2.0))
            screenshots_added = True
        except Exception:
            pass

if not screenshots_added:
    # Placeholder cards for screenshots
    placeholders = [
        ("Dashboard Overview", "Main event dashboard with KPI cards and charts"),
        ("Guest Management", "RSVP tracking with donut chart and data grid"),
        ("AI Suggestions", "Theme, menu, and venue AI-generated recommendations"),
        ("Admin Panel", "User moderation and platform statistics"),
    ]
    for i, (title, desc) in enumerate(placeholders):
        col = i % 2
        row = i // 2
        x = Inches(1.0 + col * 6.0)
        y = Inches(3.0 + row * 2.2)
        card = add_shape_rect(slide, x, y, Inches(5.5), Inches(1.9), CARD_BG, BORDER_COLOR, 0.03)
        add_text_box(slide, x + Inches(1.5), y + Inches(0.5), Inches(3.0), Inches(0.4),
                     f"[ {title} ]", font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.8), y + Inches(1.0), Inches(4.0), Inches(0.4),
                     desc, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_title(slide, "CONCLUSION", "Summary & Future Scope", 12)

# Conclusion Card
conc_card = add_shape_rect(slide, Inches(0.8), Inches(3.0), Inches(6.0), Inches(4.0), CARD_BG, BORDER_COLOR, 0.03)
add_text_box(slide, Inches(1.1), Inches(3.1), Inches(3.0), Inches(0.35),
             "✅  Key Achievements", font_size=18, color=ACCENT_GREEN, bold=True)
add_accent_line(slide, Inches(1.1), Inches(3.5), Inches(1.5), ACCENT_GREEN)

conclusions = [
    "Successfully developed a full-stack AI-powered event planning platform.",
    "Integrated multi-provider AI orchestration (Gemini, Groq, Grok) with fallback templates.",
    "Achieved 100% pass rate across 32 functional test cases.",
    "Implemented dual-mode database engine for zero-configuration execution.",
    "Demonstrated practical Generative AI embedding into CRUD applications.",
]
for i, c in enumerate(conclusions):
    y = Inches(3.7 + i * 0.55)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.08), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_GREEN
    dot.line.fill.background()
    add_text_box(slide, Inches(1.45), y - Inches(0.02), Inches(5.1), Inches(0.55),
                 c, font_size=13, color=LIGHT_GRAY)

# Future Scope Card
future_card = add_shape_rect(slide, Inches(7.1), Inches(3.0), Inches(5.5), Inches(4.0), CARD_BG, BORDER_COLOR, 0.03)
add_text_box(slide, Inches(7.4), Inches(3.1), Inches(3.0), Inches(0.35),
             "🚀  Future Scope", font_size=18, color=ACCENT_ORANGE, bold=True)
add_accent_line(slide, Inches(7.4), Inches(3.5), Inches(1.5), ACCENT_ORANGE)

future_items = [
    "Payment gateway integration (Stripe / Razorpay)",
    "Automated digital invitations via Email / SMS / WhatsApp",
    "Multi-user collaborative event workspaces using WebSockets",
    "Local vendor marketplace with reviews and direct booking",
    "Mobile application using React Native",
]
for i, f in enumerate(future_items):
    y = Inches(3.7 + i * 0.55)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), y + Inches(0.08), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_ORANGE
    dot.line.fill.background()
    add_text_box(slide, Inches(7.75), y - Inches(0.02), Inches(4.6), Inches(0.55),
                 f, font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 13 — THANK YOU
# ═══════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)

# Top bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_PURPLE
shape.line.fill.background()

# Main card
card = add_shape_rect(slide, Inches(2.5), Inches(1.5), Inches(8.333), Inches(4.5), CARD_BG, BORDER_COLOR, 0.03)

# Thank you text
add_text_box(slide, Inches(2.5), Inches(2.0), Inches(8.333), Inches(1.0),
             "THANK YOU", font_size=54, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Accent line
add_accent_line(slide, Inches(5.5), Inches(3.1), Inches(2.333), ACCENT_PURPLE)

# Subtitle
add_text_box(slide, Inches(2.5), Inches(3.4), Inches(8.333), Inches(0.5),
             "Questions & Answers Session", font_size=22, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

# Project name
add_text_box(slide, Inches(2.5), Inches(4.3), Inches(8.333), Inches(0.5),
             "JAGAH — AI Event Planner", font_size=20, color=ACCENT_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_text_box(slide, Inches(2.5), Inches(4.8), Inches(8.333), Inches(0.4),
             "Major Project Submission", font_size=14, color=MID_GRAY,
             alignment=PP_ALIGN.CENTER, font_name="Calibri")

add_bottom_bar(slide)
add_slide_number(slide, 13)

# ─── Save ────────────────────────────────────────────────────────

output_path = r"c:\Users\H\Desktop\event\event1\AI_Event_Planner_Presentation.pptx"
prs.save(output_path)
print(f"\n✅ Presentation saved successfully to:\n   {output_path}")
print(f"   Total slides: {len(prs.slides)}")
